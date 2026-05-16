"""Build the PCOS Pathfinder pptx deck.

Mirrors the 20-slide HTML deck in `presentation/pcos_pathfinder_deck.html`
plus 5 biology-context appendix slides using the matplotlib schematics in
`outputs/figures/biology/`.

Output: presentation/pcos_pathfinder_deck.pptx (16:9, 25 slides total).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

# ---------------------------------------------------------------------------
# Palette mirrors the HTML deck CSS variables.
# ---------------------------------------------------------------------------
NAVY     = RGBColor(0x1f, 0x3a, 0x5f)
INK      = RGBColor(0x14, 0x21, 0x3d)
MUTED    = RGBColor(0x5f, 0x6c, 0x7b)
LINE     = RGBColor(0xd9, 0xe1, 0xe8)
PAPER    = RGBColor(0xfb, 0xfc, 0xfd)
WHITE    = RGBColor(0xff, 0xff, 0xff)
TEAL     = RGBColor(0x00, 0x7c, 0x89)
CORAL    = RGBColor(0xd9, 0x5d, 0x39)
GOLD     = RGBColor(0xc9, 0x89, 0x10)
GREEN    = RGBColor(0x2f, 0x85, 0x5a)
LAV      = RGBColor(0x6b, 0x5b, 0x95)
SOFT_TEAL  = RGBColor(0xe7, 0xf4, 0xf5)
SOFT_CORAL = RGBColor(0xff, 0xf0, 0xeb)
SOFT_GOLD  = RGBColor(0xff, 0xf8, 0xe6)
SOFT_LAV   = RGBColor(0xef, 0xea, 0xf6)
CARD_BG    = RGBColor(0xff, 0xff, 0xff)
CODE_BG    = RGBColor(0x0f, 0x17, 0x2a)   # dark slate
CODE_FG    = RGBColor(0xe2, 0xe8, 0xf0)
CODE_KW    = RGBColor(0xc0, 0x84, 0xfc)   # purple
CODE_FN    = RGBColor(0x60, 0xa5, 0xfa)   # blue
CODE_STR   = RGBColor(0xf4, 0xa3, 0x80)   # coral
CODE_NUM   = RGBColor(0xfa, 0xcc, 0x15)   # gold
CODE_CMT   = RGBColor(0x94, 0xa3, 0xb8)   # cool grey

# Slide geometry: 13.333" x 7.5" (16:9 widescreen).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD = Inches(0.55)
CONTENT_W = SLIDE_W - 2 * PAD
TOP = Inches(0.45)

# Heights of the standard top zone (eyebrow + title).
EYEBROW_H = Inches(0.35)
TITLE_TOP = Inches(0.85)
TITLE_H = Inches(1.20)

REPO_ROOT = Path(__file__).resolve().parent.parent
FIG = REPO_ROOT / "outputs" / "figures"
BIO_FIG = FIG / "biology"
OUT_PPTX = REPO_ROOT / "presentation" / "pcos_pathfinder_deck.pptx"


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_line(shape):
    shape.line.fill.background()


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_text_box(slide, x, y, w, h, *, anchor="top", margin=Inches(0.08)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    return tb, tf


def write_runs(p, segments, *, default_size=14, default_color=INK,
               default_bold=False, default_font="Inter", default_italic=False):
    """Write a paragraph from a list of (text, **opts) segments."""
    for i, seg in enumerate(segments):
        if isinstance(seg, str):
            text, opts = seg, {}
        else:
            text, opts = seg
        run = p.add_run() if i > 0 or p.text else p.add_run()
        if i == 0 and not p.text:
            run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        f = run.font
        f.name = opts.get("font", default_font)
        f.size = Pt(opts.get("size", default_size))
        f.bold = opts.get("bold", default_bold)
        f.italic = opts.get("italic", default_italic)
        f.color.rgb = opts.get("color", default_color)


def add_paragraph(tf, text, *, size=14, color=INK, bold=False, italic=False,
                  font="Inter", align="left", space_before=0, space_after=4,
                  first=False):
    if first and tf.paragraphs and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return p


def add_runs_paragraph(tf, runs_spec, *, align="left", space_before=0,
                       space_after=4, first=False):
    """Add a paragraph composed of typed runs."""
    if first and tf.paragraphs and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    for spec in runs_spec:
        text = spec.get("text", "")
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = spec.get("font", "Inter")
        f.size = Pt(spec.get("size", 14))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.color.rgb = spec.get("color", INK)
    return p


# ---------------------------------------------------------------------------
# Slide-level helpers
# ---------------------------------------------------------------------------
def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    # Paper background.
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, PAPER)
    set_no_line(bg)
    # Move background to back.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def add_accent_strip(slide):
    """4-segment gradient strip across the bottom (teal/coral/gold/lav)."""
    h = Inches(0.10)
    y = SLIDE_H - h
    seg_w = SLIDE_W / 4
    for i, c in enumerate([TEAL, CORAL, GOLD, LAV]):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, seg_w * i, y, seg_w, h)
        set_fill(r, c)
        set_no_line(r)


def add_eyebrow(slide, text, *, color=TEAL, x=None, y=None):
    if x is None: x = PAD
    if y is None: y = TOP
    tb, tf = add_text_box(slide, x, y, Inches(11), EYEBROW_H, margin=Inches(0.0))
    add_paragraph(tf, text.upper(), size=11, color=color, bold=True,
                  space_after=0, first=True)
    return tb


def add_title(slide, text, *, color=NAVY, size=32, y=None, h=None, w=None):
    if y is None: y = TITLE_TOP
    if h is None: h = TITLE_H
    if w is None: w = CONTENT_W
    tb, tf = add_text_box(slide, PAD, y, w, h, margin=Inches(0.0))
    add_paragraph(tf, text, size=size, color=color, bold=True,
                  space_after=0, first=True)
    return tb


def add_footer(slide, slide_num, total, time_label=None):
    # Slide number only, bottom right. `time_label` accepted for backwards
    # compatibility with existing call sites but intentionally not rendered.
    tb, tf = add_text_box(slide, SLIDE_W - Inches(3.0), SLIDE_H - Inches(0.50),
                          Inches(2.7), Inches(0.30), margin=Inches(0.0))
    add_paragraph(tf, f"{slide_num} / {total}",
                  size=10, color=MUTED, align="right", first=True)


def add_card(slide, x, y, w, h, *, fill=CARD_BG, accent=None, accent_side="top",
             radius_in=0.18):
    """Rounded white card with optional accent strip."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # Reduce the corner radius for a flatter modern look.
    card.adjustments[0] = 0.05
    set_fill(card, fill)
    set_line(card, LINE, width_pt=0.75)
    if accent:
        # Use a thin accent rectangle on top edge.
        if accent_side == "top":
            ar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w,
                                        Inches(0.07))
            set_fill(ar, accent)
            set_no_line(ar)
        elif accent_side == "left":
            ar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                        Inches(0.10), h)
            set_fill(ar, accent)
            set_no_line(ar)
    return card


def add_card_content(slide, x, y, w, h, title, body, *,
                     accent=None, accent_side="top",
                     body_size=12.5, title_size=15.5,
                     fill=CARD_BG, title_color=NAVY, body_color=MUTED):
    card = add_card(slide, x, y, w, h, fill=fill, accent=accent,
                    accent_side=accent_side)
    tb, tf = add_text_box(slide, x + Inches(0.18), y + Inches(0.20),
                          w - Inches(0.36), h - Inches(0.30))
    add_paragraph(tf, title, size=title_size, color=title_color, bold=True,
                  space_after=4, first=True)
    if isinstance(body, str):
        add_paragraph(tf, body, size=body_size, color=body_color, space_after=0)
    else:
        for line in body:
            add_paragraph(tf, line, size=body_size, color=body_color, space_after=2)
    return card


def add_metric_card(slide, x, y, w, h, value, label, *, accent=TEAL):
    card = add_card(slide, x, y, w, h, accent=accent, fill=CARD_BG)
    tb, tf = add_text_box(slide, x, y + Inches(0.20),
                          w, h - Inches(0.30), margin=Inches(0.0))
    add_paragraph(tf, value, size=36, color=accent, bold=True,
                  align="center", space_after=4, first=True)
    add_paragraph(tf, label, size=11, color=MUTED, align="center",
                  space_after=0)
    return card


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        # Render a placeholder if the asset is missing.
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                      w or Inches(4), h or Inches(2.5))
        set_fill(rect, SOFT_TEAL)
        set_line(rect, LINE, width_pt=0.5)
        tb, tf = add_text_box(slide, x, y, w or Inches(4), h or Inches(2.5),
                              anchor="middle", margin=Inches(0))
        add_paragraph(tf, f"[missing: {Path(path).name}]", size=11,
                      color=MUTED, align="center", first=True)
        return rect
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    return pic


def add_image_centered(slide, path, x, y, w, h):
    """Embed an image while preserving aspect ratio and centering inside box."""
    if not Path(path).exists():
        return add_image(slide, path, x, y, w=w, h=h)
    pic = slide.shapes.add_picture(str(path), x, y)
    iw, ih = pic.width, pic.height
    scale = min(w / iw, h / ih)
    new_w = int(iw * scale)
    new_h = int(ih * scale)
    pic.width = new_w
    pic.height = new_h
    pic.left = int(x + (w - new_w) / 2)
    pic.top  = int(y + (h - new_h) / 2)
    return pic


def add_caption(slide, x, y, w, text):
    tb, tf = add_text_box(slide, x, y, w, Inches(0.30), margin=Inches(0.0))
    add_paragraph(tf, text, size=10, color=MUTED, italic=True,
                  align="center", first=True)


# ---------------------------------------------------------------------------
# Code/output panels for demo slides
# ---------------------------------------------------------------------------
@dataclass
class CodeLine:
    parts: list[tuple[str, str]]  # (token-class, text)


def code_card(slide, x, y, w, h, label, lines, *, label_color=TEAL):
    """Dark code panel with monospace text and simple syntax-colouring tokens."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.adjustments[0] = 0.04
    set_fill(panel, CODE_BG)
    set_line(panel, RGBColor(0x1e, 0x29, 0x3b), width_pt=0.5)

    # Label band.
    lab_h = Inches(0.35)
    lab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, lab_h)
    set_fill(lab, RGBColor(0x1e, 0x29, 0x3b))
    set_no_line(lab)
    tb, tf = add_text_box(slide, x + Inches(0.18), y, w - Inches(0.30),
                          lab_h, anchor="middle", margin=Inches(0.0))
    add_paragraph(tf, label, size=10, color=RGBColor(0x9b, 0xb1, 0xd0),
                  font="Consolas", bold=True, first=True)

    # Code text.
    code_y = y + lab_h + Inches(0.05)
    code_h = h - lab_h - Inches(0.10)
    tb, tf = add_text_box(slide, x + Inches(0.20), code_y, w - Inches(0.30),
                          code_h, margin=Inches(0.05))
    token_colors = {
        "fg": CODE_FG, "k": CODE_KW, "f": CODE_FN, "s": CODE_STR,
        "n": CODE_NUM, "c": CODE_CMT,
    }
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        if isinstance(line, str):
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = "Consolas"
            run.font.size = Pt(10.5)
            run.font.color.rgb = CODE_FG
            continue
        for tok, text in line.parts:
            run = p.add_run()
            run.text = text
            run.font.name = "Consolas"
            run.font.size = Pt(10.5)
            run.font.color.rgb = token_colors.get(tok, CODE_FG)
    return panel


def output_card(slide, x, y, w, h, label, lines):
    """Light output panel with monospace text and highlight tokens."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.adjustments[0] = 0.04
    set_fill(panel, RGBColor(0xf8, 0xfa, 0xfc))
    set_line(panel, LINE, width_pt=0.6)

    lab_h = Inches(0.35)
    lab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, lab_h)
    set_fill(lab, SOFT_TEAL)
    set_no_line(lab)
    tb, tf = add_text_box(slide, x + Inches(0.18), y, w - Inches(0.30),
                          lab_h, anchor="middle", margin=Inches(0.0))
    add_paragraph(tf, label, size=10, color=TEAL, font="Consolas",
                  bold=True, first=True)

    out_y = y + lab_h + Inches(0.05)
    out_h = h - lab_h - Inches(0.10)
    tb, tf = add_text_box(slide, x + Inches(0.20), out_y, w - Inches(0.30),
                          out_h, margin=Inches(0.05))
    token_colors = {
        "fg": INK, "mu": MUTED, "ok": GREEN, "hi": CORAL,
    }
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        if isinstance(line, str):
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = "Consolas"
            run.font.size = Pt(10.5)
            run.font.color.rgb = INK
            continue
        for tok, text in line.parts:
            run = p.add_run()
            run.text = text
            run.font.name = "Consolas"
            run.font.size = Pt(10.5)
            run.font.color.rgb = token_colors.get(tok, INK)
    return panel


def takeaway(slide, x, y, w, h, text):
    """Coral takeaway band at the bottom of a demo slide."""
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    band.adjustments[0] = 0.18
    set_fill(band, SOFT_CORAL)
    set_line(band, CORAL, width_pt=1.0)
    tb, tf = add_text_box(slide, x + Inches(0.30), y, w - Inches(0.55), h,
                          anchor="middle", margin=Inches(0.0))
    add_runs_paragraph(tf, [
        {"text": "Takeaway   ", "size": 11, "color": CORAL, "bold": True},
        {"text": text, "size": 12.5, "color": INK},
    ], first=True)


def line(parts):
    return CodeLine(parts=parts)


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text
    for p in notes.paragraphs:
        for r in p.runs:
            r.font.name = "Inter"
            r.font.size = Pt(13)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_title(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Clinical decision support  ·  Biohackathon 2026 Women's Health")

    # Big title.
    tb, tf = add_text_box(slide, PAD, Inches(1.10), Inches(7.5), Inches(2.0),
                          margin=Inches(0))
    add_paragraph(tf, "PCOS Pathfinder", size=64, color=NAVY, bold=True,
                  first=True)

    # Subtitle.
    tb, tf = add_text_box(slide, PAD, Inches(3.25), Inches(7.4), Inches(1.6),
                          margin=Inches(0))
    add_paragraph(tf, "A tiered, explainable clinical decision-support prototype "
                  "for earlier PCOS detection and safer differential workup.",
                  size=20, color=RGBColor(0x33, 0x41, 0x55), first=True)

    # Metric strip.
    card_w = Inches(2.35)
    card_h = Inches(1.40)
    y = Inches(5.10)
    gap = Inches(0.20)
    x = PAD
    for value, label, color in [
        ("0.953", "Enhanced AUC",        TEAL),
        ("0.886", "Sensitivity",         CORAL),
        ("13",    "Reproducible notebooks", GOLD),
    ]:
        add_metric_card(slide, x, y, card_w, card_h, value, label, accent=color)
        x += card_w + gap

    # ROC figure on the right.
    fig_x = Inches(8.40)
    fig_y = Inches(1.00)
    fig_w = Inches(4.45)
    fig_h = Inches(4.40)
    add_image_centered(slide, FIG / "enhanced_roc_curve.png",
                       fig_x, fig_y, fig_w, fig_h)
    add_caption(slide, fig_x, fig_y + fig_h + Inches(0.05), fig_w,
                "Held-out test split, n=136. Thresholds selected from training "
                "out-of-fold probabilities.")

    add_footer(slide, n, total, "0:00 – 0:25")
    set_notes(slide,
        "Good morning. We built PCOS Pathfinder: a practical, explainable "
        "decision-support workflow for earlier PCOS detection. Our thesis is "
        "simple — clinical AI should not just predict a label. It should help "
        "clinicians decide who needs assessment, explain why, and avoid "
        "missing overlapping conditions. On the held-out split, the enhanced "
        "model reaches 0.953 ROC-AUC with 0.886 sensitivity, and everything is "
        "reproducible from thirteen notebooks.")


def build_problem(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "The diagnostic gap")
    add_title(slide, "PCOS is common, heterogeneous, and routinely missed.", size=28)

    # 3 stat pills — concrete prevalence/delay/under-diagnosis numbers.
    stat_w = Inches(2.45)
    stat_h = Inches(1.30)
    stat_y = Inches(2.05)
    sx = PAD
    for value, label, color in [
        ("8–13%",  "of reproductive-age women affected\n(Teede et al., Fertil Steril 2023)", TEAL),
        ("4+ yrs", "average diagnostic delay\nacross studies",                              CORAL),
        ("~70%",   "estimated undiagnosed globally\n(Teede 2023; March 2010)",              GOLD),
    ]:
        card = add_card(slide, sx, stat_y, stat_w, stat_h, accent=color)
        tb, tf = add_text_box(slide, sx, stat_y + Inches(0.18),
                              stat_w, stat_h - Inches(0.30), margin=Inches(0.10))
        add_paragraph(tf, value, size=26, color=color, bold=True,
                      align="center", space_after=4, first=True)
        add_paragraph(tf, label, size=10, color=MUTED, align="center")
        sx += stat_w + Inches(0.10)

    # 3 cards on why detection is hard.
    cards = [
        ("Complex presentation",
         "Acne, weight change, irregular bleeding, infertility, metabolic risk, "
         "and ultrasound findings can each be the first signal.", TEAL),
        ("Delayed pathway",
         "Primary care sees the earliest signals; labs and ultrasound often "
         "arrive only after specialist referral.", CORAL),
        ("Overlap risk",
         "Pelvic pain and infertility can point toward adjacent gynecologic "
         "conditions — the system must reduce tunnel vision.", GOLD),
    ]
    card_w = Inches(2.45)
    card_h = Inches(2.30)
    y = Inches(3.55)
    x = PAD
    for title, body, color in cards:
        add_card_content(slide, x, y, card_w, card_h, title, body,
                         accent=color, body_size=12)
        x += card_w + Inches(0.10)

    # Right-hand figure.
    fig_x = Inches(8.40)
    fig_y = Inches(2.00)
    fig_w = Inches(4.45)
    fig_h = Inches(4.30)
    add_image_centered(slide, FIG / "pcos_symptom_rates.png",
                       fig_x, fig_y, fig_w, fig_h)
    add_caption(slide, fig_x, fig_y + fig_h + Inches(0.00), fig_w,
                "Source cohort symptom differences ground the clinical story.")

    add_footer(slide, n, total, "0:25 – 1:00")
    set_notes(slide,
        "The challenge asks us to improve diagnostic accuracy in women's "
        "health. PCOS is a perfect core case because it is common, "
        "heterogeneous, and often delayed. A patient may first present with "
        "acne, hair growth, irregular bleeding, weight changes, or fertility "
        "concerns. Those signals are clinically meaningful, but they are often "
        "spread across visits and specialists. Our design goal is not just a "
        "high AUC — it is a workflow that fits the diagnostic pathway.")


def build_workflow(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Tiered workflow")
    add_title(slide, "A tiered pathway that matches real clinical escalation.",
              size=30)

    # 4 numbered steps.
    steps = [
        ("1", "Frontline screening",
         "Symptoms and basic vitals for primary care, telehealth, or outreach.", TEAL),
        ("2", "Enhanced support",
         "Add labs and ultrasound when the clinician has more evidence.", CORAL),
        ("3", "Checklist",
         "Rotterdam-style completeness, missing-testosterone caveat, metabolic "
         "prompt, lean-PCOS warning.", GOLD),
        ("4", "Differential prompt",
         "Endometriosis-overlap nudge using synthetic data, framed only as "
         "workflow support.", LAV),
    ]
    step_w = Inches(2.95)
    step_h = Inches(3.10)
    y = Inches(2.40)
    gap = Inches(0.15)
    x = PAD
    for num, t, body, color in steps:
        card = add_card(slide, x, y, step_w, step_h, accent=color)
        # Big number badge.
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25),
                                       y + Inches(0.25), Inches(0.75),
                                       Inches(0.75))
        set_fill(badge, color)
        set_no_line(badge)
        tb, tf = add_text_box(slide, x + Inches(0.25), y + Inches(0.25),
                              Inches(0.75), Inches(0.75), anchor="middle",
                              margin=Inches(0))
        add_paragraph(tf, num, size=22, color=WHITE, bold=True,
                      align="center", first=True)
        # Title + body.
        tb, tf = add_text_box(slide, x + Inches(0.25), y + Inches(1.15),
                              step_w - Inches(0.50), step_h - Inches(1.30))
        add_paragraph(tf, t, size=17, color=NAVY, bold=True, space_after=8,
                      first=True)
        add_paragraph(tf, body, size=12.5, color=MUTED)
        x += step_w + gap

    # Bottom callout.
    cb_w = CONTENT_W
    cb_h = Inches(0.85)
    cb_y = Inches(5.85)
    card = add_card(slide, PAD, cb_y, cb_w, cb_h, fill=SOFT_TEAL,
                    accent=TEAL, accent_side="left")
    tb, tf = add_text_box(slide, PAD + Inches(0.30), cb_y, cb_w - Inches(0.40),
                          cb_h, anchor="middle", margin=Inches(0))
    add_runs_paragraph(tf, [
        {"text": "Core idea  ", "size": 13, "color": TEAL, "bold": True},
        {"text": "Give clinicians a sequence of next-best actions, not an "
                 "opaque yes/no diagnosis.", "size": 14, "color": INK},
    ], first=True)

    add_footer(slide, n, total, "3:55 – 4:30")
    set_notes(slide,
        "The product is tiered. First, a frontline model uses low-friction "
        "data. Second, an enhanced model uses labs and ultrasound. Third, the "
        "app shows a diagnostic-completeness checklist so the clinician sees "
        "what evidence is present and what is missing. Fourth, the overlap "
        "prompt warns when the symptom pattern deserves broader gynecologic "
        "workup. That is why this is decision support, not automated diagnosis.")


def build_data(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Data foundation")
    add_title(slide, "Data choices were conservative by design.", size=30)

    # Left: dataset table. PCOS row surfaces single-source Kerala origin.
    tbl_x = PAD
    tbl_y = Inches(2.20)
    tbl_w = Inches(7.20)
    headers = ["Dataset", "Use", "Claim boundary"]
    rows = [
        ["PCOS clinical, n=541",
         "Main modelling foundation",
         "10 hospitals, Kerala, India · external validation required"],
        ["Endometriosis, n=10,000",
         "Symptom-overlap prompt",
         "Synthetic; not diagnostic validation"],
        ["Single-cell archives",
         "Gene-list rationale",
         "Feature-list peek only; no expression analysis"],
    ]
    _add_table(slide, tbl_x, tbl_y, tbl_w, headers, rows,
               col_widths=[Inches(2.30), Inches(2.30), Inches(2.60)],
               header_fill=NAVY, header_color=WHITE)

    # Right: 4 cards stacked.
    cards = [
        ("Cycle-length data-quality fix",
         "Source column labelled \"cycle length\" actually ranges 0–12 days — "
         "this is bleeding duration, not cycle interval. App relabels it to "
         "avoid out-of-distribution inputs.", CORAL),
        ("Blood-group caveat",
         "Yang et al. 2025 supports phenotype/severity context, not ordinal "
         "model input or standalone diagnosis.", TEAL),
        ("Excluded from screening",
         "Fast food and exercise proxies, to avoid stigma and weak "
         "generalization.", GOLD),
        ("Coercions logged",
         "Every non-numeric conversion is exported to the audit files.", LAV),
    ]
    cx = Inches(7.95)
    cy = Inches(2.20)
    cw = Inches(4.85)
    ch = Inches(1.14)
    gap = Inches(0.12)
    for title, body, color in cards:
        add_card_content(slide, cx, cy, cw, ch, title, body, accent=color,
                         accent_side="left", body_size=11.5, title_size=13.5)
        cy += ch + gap

    add_footer(slide, n, total, "4:30 – 5:05")
    set_notes(slide,
        "The data boundaries are as important as the model. We use the PCOS "
        "clinical data as the foundation. The endometriosis dataset is "
        "synthetic, so we only use it as a workflow prompt. The single-cell "
        "data supports biological rationale only. We also fixed a critical "
        "data-quality issue: the source says cycle length, but the range is 0 "
        "to 12 days, so it is actually bleeding duration. That prevents "
        "out-of-distribution inputs like 28 days. Blood group is handled as a "
        "categorical caveat per Yang et al. 2025.")


def _add_table(slide, x, y, w, headers, rows, *, col_widths=None,
               header_fill=NAVY, header_color=WHITE,
               row_h=Inches(0.45), header_h=Inches(0.42),
               highlight_rows=None):
    n_cols = len(headers)
    if col_widths is None:
        col_w = w / n_cols
        col_widths = [col_w] * n_cols
    table = slide.shapes.add_table(len(rows) + 1, n_cols, x, y, w,
                                   header_h + row_h * len(rows)).table
    # Column widths.
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    # Row heights.
    table.rows[0].height = header_h
    for i in range(len(rows)):
        table.rows[i + 1].height = row_h
    # Header.
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = h
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Inter"
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = header_color
    # Body.
    highlight_rows = highlight_rows or {}
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 \
                else RGBColor(0xf8, 0xfa, 0xfc)
            tf = cell.text_frame
            tf.margin_left = Inches(0.10)
            tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = str(val)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = "Inter"
                    r.font.size = Pt(11.5)
                    r.font.color.rgb = INK
                    if (ri, ci) in highlight_rows:
                        r.font.color.rgb = highlight_rows[(ri, ci)]
                        r.font.bold = True
    return table


def build_demo1(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Demo 1  ·  notebook 01  ·  appendix", color=CORAL)
    add_title(slide,
              "Every silently coerced cell is logged. Every dropped column is justified.",
              size=24, y=Inches(0.78), h=Inches(0.95))

    code_lines = [
        line([("c", "# Drop identifiers / sensitive fields; withhold")]),
        line([("c", "# blood_group until categorical encoding is validated.")]),
        line([("k", "def "), ("f", "clean_pcos_dataframe"), ("fg", "(df):")]),
        line([("fg", "    df = df.copy()")]),
        line([("fg", "    coercion_report = {}")]),
        line([("k", "    for"), ("fg", " col "), ("k", "in"), ("fg", " df.columns:")]),
        line([("fg", "        coerced = pd."), ("f", "to_numeric"),
              ("fg", "(df[col], errors="), ("s", '"coerce"'), ("fg", ")")]),
        line([("fg", "        new_nan = coerced.isna()."),
              ("f", "sum"), ("fg", "() - df[col].isna()."),
              ("f", "sum"), ("fg", "()")]),
        line([("k", "        if"), ("fg", " new_nan > "), ("n", "0"), ("fg", ":")]),
        line([("fg", "            offenders = df.loc[df[col]."),
              ("f", "notna"), ("fg", "() & coerced."),
              ("f", "isna"), ("fg", "(), col]")]),
        line([("fg", "            coercion_report[col] = {")]),
        line([("fg", "                "), ("s", '"coerced_to_nan"'),
              ("fg", ": "), ("f", "int"), ("fg", "(new_nan),")]),
        line([("fg", "                "), ("s", '"examples"'),
              ("fg", ": offenders."), ("f", "astype"),
              ("fg", "("), ("f", "str"),
              ("fg", ")."), ("f", "unique"),
              ("fg", "()."), ("f", "tolist"),
              ("fg", "()[:"), ("n", "5"), ("fg", "],")]),
        line([("fg", "            }")]),
        line([("fg", "        df[col] = coerced")]),
        line([("fg", "    df = df.drop(columns=[c "), ("k", "for"),
              ("fg", " c "), ("k", "in"), ("fg", " DROP_COLUMNS "),
              ("k", "if"), ("fg", " c "), ("k", "in"), ("fg", " df.columns])")]),
        line([("fg", "    df.attrs["), ("s", '"coercion_report"'),
              ("fg", "] = coercion_report")]),
        line([("k", "    return"), ("fg", " df")]),
        "",
        line([("fg", "DROP_COLUMNS = ["), ("s", '"sl_no"'),
              ("fg", ", "), ("s", '"patient_file_no"'), ("fg", ",")]),
        line([("fg", "                "), ("s", '"blood_group"'),
              ("fg", ", "), ("s", '"marriage_status_yrs"'), ("fg", "]")]),
    ]
    code_card(slide, PAD, Inches(1.85), Inches(6.30), Inches(4.55),
              "scripts/create_training_notebooks.mjs  —  clean_pcos_dataframe",
              code_lines)

    out_lines = [
        line([("mu", "PCOS shape: (541, 41)")]),
        line([("mu", "Dropped columns: ['sl_no', 'patient_file_no',")]),
        line([("mu", "                  'blood_group', 'marriage_status_yrs']")]),
        line([("mu", "Non-numeric values coerced to NaN:")]),
        line([("fg", "  "), ("hi", "ii_beta_hcg_miu_ml"),
              ("fg", ":  1 cell(s); examples=['1.99.']")]),
        line([("fg", "  "), ("hi", "amh_ng_ml"),
              ("fg", ":           1 cell(s); examples=['a']")]),
        "",
        line([("mu", "pcos_y_n")]),
        line([("fg", "0    "), ("ok", "364")]),
        line([("fg", "1    "), ("ok", "177")]),
        line([("fg", "Name: count, dtype: int64")]),
    ]
    output_card(slide, Inches(7.10), Inches(1.85), Inches(5.75), Inches(4.55),
                "Output  —  outputs/metrics/pcos_coercion_report.csv",
                out_lines)

    takeaway(slide, PAD, Inches(6.55), CONTENT_W, Inches(0.70),
        "Two stray data-quality issues caught and logged before they can bias a model. "
        "ABO/Rh withheld from screening; Yang et al. 2025 supports retaining it later "
        "as categorical phenotype metadata, not as a numeric 11–18 input.")

    add_footer(slide, n, total, "appendix / Q&A")
    set_notes(slide,
        "This is the cleaning step. Every column gets numeric-coerced and "
        "every non-numeric cell that drops to NaN is recorded with an example. "
        "That is how we caught the literal string '1.99.' in the second "
        "beta-HCG column and the literal letter 'a' in the AMH column. We "
        "drop identifiers and the marriage-status column. Blood group is "
        "withheld because codes 11–18 are categorical ABO/Rh labels, not an "
        "ordered diagnostic scale.")


def build_method(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Model design")
    add_title(slide, "Validation avoids the common traps.", size=30)

    # Left column: 3 cards stacked.
    cards = [
        ("No test leakage",
         "Model selection uses 5-fold CV on training data. Action thresholds "
         "come from out-of-fold training probabilities.", TEAL),
        ("One final holdout",
         "The held-out test split is touched once for final reporting.", CORAL),
        ("Recall-first thresholding",
         "Threshold sweep targets training-CV recall ≥ 0.90, then maximizes "
         "specificity, precision, and F1.", GOLD),
    ]
    cx = PAD
    cy = Inches(2.25)
    cw = Inches(5.60)
    ch = Inches(1.25)
    for t, body, color in cards:
        add_card_content(slide, cx, cy, cw, ch, t, body, accent=color,
                         accent_side="left", body_size=12, title_size=14.5)
        cy += ch + Inches(0.18)

    # Right column: model table + PR figure.
    headers = ["Model", "Feature set", "Clinical role"]
    rows = [
        ["Screening RF",     "13 basic features",   "Rule-out triage"],
        ["Enhanced RF",      "27 clinical features", "Diagnostic support"],
        ["Endometriosis LR", "6 overlap features",  "Differential prompt only"],
    ]
    _add_table(slide, Inches(7.05), Inches(2.25), Inches(5.80),
               headers, rows,
               col_widths=[Inches(1.90), Inches(2.00), Inches(1.90)])
    # PR figure.
    add_image_centered(slide, FIG / "screening_precision_recall_curve.png",
                       Inches(7.05), Inches(4.40), Inches(5.80), Inches(2.55))

    add_footer(slide, n, total, "5:05 – 5:40")
    set_notes(slide,
        "The threshold is not chosen on the test set. We select models using "
        "cross-validation on training data, then choose thresholds using "
        "out-of-fold probabilities from that training split. Only after that "
        "do we evaluate on the held-out test set. We optimize for recall — "
        "missing a PCOS case is more costly than sending someone for follow-up.")


def build_demo2(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Demo 2  ·  notebook 02  ·  appendix", color=CORAL)
    add_title(slide,
              "Action thresholds come from training cross-validation. The test set is touched once.",
              size=22, y=Inches(0.78), h=Inches(0.95))

    code_lines = [
        line([("c", "# 1. Out-of-fold probabilities on the TRAINING set only.")]),
        line([("fg", "cv = "), ("f", "StratifiedKFold"),
              ("fg", "(n_splits="), ("n", "5"),
              ("fg", ", shuffle="), ("k", "True"), ("fg", ",")]),
        line([("fg", "                     random_state=RANDOM_STATE)")]),
        line([("fg", "oof_proba = "), ("f", "cross_val_predict"), ("fg", "(")]),
        line([("fg", "    best_model, X_train, y_train,")]),
        line([("fg", "    cv=cv, method="), ("s", '"predict_proba"'),
              ("fg", ", n_jobs=-"), ("n", "1"), ("fg", ",")]),
        line([("fg", ")[:, "), ("n", "1"), ("fg", "]")]),
        "",
        line([("c", "# 2. Sweep thresholds; pick the one with recall >= 0.90")]),
        line([("c", "#    that also maximises specificity, precision, F1.")]),
        line([("fg", "chosen, table = "), ("f", "choose_high_recall_threshold"),
              ("fg", "(")]),
        line([("fg", "    y_train, oof_proba, min_recall="), ("n", "0.90"), ("fg", ",")]),
        line([("fg", ")")]),
        line([("fg", "threshold = "), ("f", "float"), ("fg", "(chosen["),
              ("s", '"threshold"'), ("fg", "])")]),
        "",
        line([("c", "# 3. Refit on FULL training data, evaluate ONCE on test.")]),
        line([("fg", "best_model."), ("f", "fit"), ("fg", "(X_train, y_train)")]),
        line([("fg", "metrics = "), ("f", "evaluate_holdout"),
              ("fg", "(best_model, X_test, y_test,")]),
        line([("fg", "                           threshold=threshold)")]),
    ]
    code_card(slide, PAD, Inches(1.85), Inches(6.30), Inches(4.55),
              "notebook 02  —  cv_oof_probabilities + threshold sweep",
              code_lines)

    out_lines = [
        line([("mu", "Best screening model: random_forest")]),
        line([("mu", "Chosen threshold (from training CV):")]),
        line([("fg", "  threshold           "), ("ok", "0.285")]),
        line([("fg", "  recall              "), ("ok", "0.901")]),
        line([("fg", "  specificity         0.683")]),
        line([("fg", "  precision           0.535")]),
        line([("fg", "  f1                  0.669")]),
        line([("fg", "  balanced_accuracy   0.792")]),
        "",
        line([("mu", "Holdout evaluation (touched ONCE):")]),
        line([("fg", "  roc_auc       "), ("ok", "0.896")]),
        line([("fg", "  recall        "), ("ok", "0.886")]),
        line([("fg", "  specificity   0.685")]),
        line([("fg", "  precision     0.574")]),
        line([("fg", "  npv           0.926")]),
        line([("fg", "  f1            0.696")]),
        line([("fg", "  f2            0.799")]),
    ]
    output_card(slide, Inches(7.10), Inches(1.85), Inches(5.75), Inches(4.55),
                "Output  —  chosen threshold + held-out test",
                out_lines)

    takeaway(slide, PAD, Inches(6.55), CONTENT_W, Inches(0.70),
        "The test set is evaluated exactly once after the model is refit on the full "
        "training data — that is why the recall and specificity numbers on the next "
        "slide are honest.")

    add_footer(slide, n, total, "appendix / Q&A")
    set_notes(slide,
        "A common failure mode in clinical ML pitches is choosing the "
        "threshold on the same test set you then report metrics from. We "
        "avoid it. The threshold is picked from out-of-fold training "
        "probabilities. The held-out test set is evaluated exactly once, "
        "after the model is refit on the full training data.")


def build_results(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Held-out performance")
    add_title(slide, "Enhanced support preserves sensitivity and greatly improves precision.",
              size=27)

    headers = ["Model", "Threshold", "AUC", "Recall", "Specificity", "PPV", "NPV", "F2"]
    rows = [
        ["Screening",            "0.285", "0.896", "0.886", "0.685", "0.574", "0.926", "0.799"],
        ["Enhanced",             "0.380", "0.953", "0.886", "0.902", "0.812", "0.943", "0.871"],
        ["Endometriosis prompt", "0.510", "0.660", "0.628", "0.618", "0.532", "—",     "0.609"],
    ]
    highlight = {
        (0, 3): GREEN,
        (1, 2): GREEN, (1, 3): GREEN, (1, 4): GREEN, (1, 5): GREEN,
        (1, 6): GREEN, (1, 7): GREEN,
    }
    _add_table(slide, PAD, Inches(2.20), CONTENT_W, headers, rows,
               col_widths=[Inches(2.60), Inches(1.30), Inches(1.30), Inches(1.30),
                           Inches(1.55), Inches(1.30), Inches(1.30), Inches(1.55)],
               row_h=Inches(0.55), header_h=Inches(0.50),
               highlight_rows=highlight)

    # 4 metric cards along the bottom.
    metrics = [
        ("39/44", "PCOS positives caught", TEAL),
        ("+21.7", "Specificity points vs screening", CORAL),
        ("0.943", "Enhanced NPV", GOLD),
        ("+0.057", "ΔAUC vs screening, 95% CI [+0.013, +0.105], p≈0.015", LAV),
    ]
    card_w = (CONTENT_W - Inches(0.30) * 3) / 4
    cy = Inches(4.85)
    ch = Inches(1.80)
    x = PAD
    for value, label, color in metrics:
        card = add_card(slide, x, cy, card_w, ch, accent=color)
        tb, tf = add_text_box(slide, x, cy + Inches(0.25),
                              card_w, ch - Inches(0.40), margin=Inches(0.10))
        add_paragraph(tf, value, size=34, color=color, bold=True,
                      align="center", space_after=8, first=True)
        add_paragraph(tf, label, size=10.5, color=MUTED, align="center")
        x += card_w + Inches(0.30)

    add_footer(slide, n, total, "5:40 – 6:20")
    set_notes(slide,
        "The screening model catches 39 of 44 PCOS-positive patients in the "
        "holdout cohort. The enhanced model keeps the same sensitivity, but "
        "raises specificity from 0.685 to 0.902 and precision from 0.574 to "
        "0.812. That is the point of the tiered workflow: screen broadly "
        "first, then use labs and ultrasound to reduce false positives.")


def build_rigor(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Uncertainty, calibration & conformal  ·  TRIPOD+AI aligned")
    add_title(slide, "The rigor layer makes the model harder to dismiss.", size=30)

    cards = [
        ("Uncertainty",
         "2000-resample bootstrap CIs; enhanced AUC advantage: +0.057, 95% CI [+0.013, +0.105].",
         TEAL),
        ("Calibration",
         "Enhanced Platt scaling improves Brier 0.093 → 0.072 and ECE 0.143 → 0.045.",
         CORAL),
        ("Conformal coverage",
         "Target 0.90; empirical holdout coverage 0.912. Empty sets are abstentions.",
         GOLD),
    ]
    cx = PAD
    cy = Inches(2.30)
    cw = Inches(5.60)
    ch = Inches(1.30)
    for t, body, color in cards:
        add_card_content(slide, cx, cy, cw, ch, t, body, accent=color,
                         accent_side="left", body_size=12, title_size=14.5)
        cy += ch + Inches(0.20)

    # Reporting-standard tagline below the cards.
    tb, tf = add_text_box(slide, PAD, Inches(6.40), Inches(5.60), Inches(0.55),
                          margin=Inches(0))
    add_runs_paragraph(tf, [
        {"text": "Reporting standard  ", "size": 11, "color": NAVY, "bold": True},
        {"text": "Collins et al., BMJ 2024 (TRIPOD+AI); Van Calster et al., BMC Med 2019; "
                 "Angelopoulos & Bates 2023.",
         "size": 10.5, "color": MUTED, "italic": True},
    ], first=True)

    # Right: calibration figure.
    add_image_centered(slide, FIG / "calibration_combined.png",
                       Inches(7.05), Inches(2.30), Inches(5.80), Inches(4.30))
    add_caption(slide, Inches(7.05), Inches(6.65), Inches(5.80),
                "Notebook 08: probability quality before and after recalibration.")

    add_footer(slide, n, total, "6:20 – 7:05")
    set_notes(slide,
        "Point estimates alone are not enough for a clinical prediction "
        "pitch. We added bootstrap confidence intervals, calibration, decision "
        "curve analysis, conformal prediction, subgroup audits, TabPFN "
        "benchmarking, and a Rotterdam comparison. The combined layer is what "
        "makes the model harder to dismiss in peer review.")


def build_demo3(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Demo 3  ·  notebook 07  ·  appendix", color=CORAL)
    add_title(slide,
              "Paired bootstrap (2000 resamples) gives an AUC delta with a real interval.",
              size=22, y=Inches(0.78), h=Inches(0.95))

    code_lines = [
        line([("fg", "rng = np."), ("f", "default_rng"),
              ("fg", "(seed="), ("n", "42"), ("fg", ")")]),
        line([("fg", "n = "), ("f", "len"), ("fg", "(y_test)")]),
        line([("fg", "deltas = []")]),
        line([("k", "for"), ("fg", " _ "), ("k", "in"),
              ("fg", " "), ("f", "range"), ("fg", "("), ("n", "2000"), ("fg", "):")]),
        line([("fg", "    idx = rng."), ("f", "choice"),
              ("fg", "(n, n, replace="), ("k", "True"), ("fg", ")")]),
        line([("fg", "    yb  = y_test.values[idx]")]),
        line([("k", "    if"), ("fg", " yb."), ("f", "sum"),
              ("fg", "() == "), ("n", "0"),
              ("fg", " "), ("k", "or"), ("fg", " yb."),
              ("f", "sum"), ("fg", "() == "),
              ("f", "len"), ("fg", "(yb):")]),
        line([("k", "        continue"),
              ("c", "  # AUC undefined on single-class draw")]),
        line([("fg", "    auc_s = "), ("f", "roc_auc_score"),
              ("fg", "(yb, screening_proba[idx])")]),
        line([("fg", "    auc_e = "), ("f", "roc_auc_score"),
              ("fg", "(yb, enhanced_proba[idx])")]),
        line([("fg", "    deltas."), ("f", "append"), ("fg", "(auc_e - auc_s)")]),
        "",
        line([("fg", "deltas = np."), ("f", "asarray"), ("fg", "(deltas)")]),
        line([("fg", "point = deltas."), ("f", "mean"), ("fg", "()")]),
        line([("fg", "ci    = np."), ("f", "percentile"),
              ("fg", "(deltas, ["), ("n", "2.5"), ("fg", ", "),
              ("n", "97.5"), ("fg", "])")]),
        line([("fg", "p_val = "), ("n", "2"), ("fg", " * "),
              ("f", "min"), ("fg", "((deltas <= "),
              ("n", "0"), ("fg", ")."),
              ("f", "mean"), ("fg", "(),")]),
        line([("fg", "                  (deltas >= "), ("n", "0"),
              ("fg", ")."), ("f", "mean"), ("fg", "())")]),
    ]
    code_card(slide, PAD, Inches(1.85), Inches(6.30), Inches(4.55),
              "notebook 07  —  paired bootstrap loop", code_lines)

    out_lines = [
        line([("fg", "{")]),
        line([("fg", "  "), ("mu", '"screening"'), ("fg", ": {"),
              ("mu", '"roc_auc"'), ("fg", ": {")]),
        line([("fg", "    "), ("mu", '"point"'), ("fg", ": 0.8955,")]),
        line([("fg", "    "), ("mu", '"ci_lo"'),
              ("fg", ": 0.8300, "), ("mu", '"ci_hi"'),
              ("fg", ": 0.9515")]),
        line([("fg", "  }},")]),
        line([("fg", "  "), ("mu", '"enhanced"'), ("fg", ": {"),
              ("mu", '"roc_auc"'), ("fg", ": {")]),
        line([("fg", "    "), ("mu", '"point"'), ("fg", ":  "),
              ("ok", "0.9528"), ("fg", ",")]),
        line([("fg", "    "), ("mu", '"ci_lo"'), ("fg", ":  "),
              ("ok", "0.9114"), ("fg", ", "),
              ("mu", '"ci_hi"'), ("fg", ":  "), ("ok", "0.9852")]),
        line([("fg", "  }},")]),
        line([("fg", "  "), ("mu", '"auc_difference"'), ("fg", ": {")]),
        line([("fg", "    "), ("mu", '"point"'),
              ("fg", ":               "), ("hi", "+0.0573")]),
        line([("fg", "    "), ("mu", '"ci_lo"'),
              ("fg", ":               "), ("hi", "+0.0130")]),
        line([("fg", "    "), ("mu", '"ci_hi"'),
              ("fg", ":               "), ("hi", "+0.1051")]),
        line([("fg", "    "), ("mu", '"p_value_two_sided"'),
              ("fg", ":  "), ("hi", "~0.015")]),
        line([("fg", "  }")]),
        line([("fg", "}")]),
    ]
    output_card(slide, Inches(7.10), Inches(1.85), Inches(5.75), Inches(4.55),
                "Output  —  outputs/metrics/holdout_bootstrap_cis.json",
                out_lines)

    takeaway(slide, PAD, Inches(6.55), CONTENT_W, Inches(0.70),
        "The enhanced model's AUC advantage is +0.057 with a 95% CI that excludes zero — "
        "the kind of statement that survives peer review.")

    add_footer(slide, n, total, "appendix / Q&A")
    set_notes(slide,
        "We ran a paired bootstrap: the same 2000 resampled index vectors for "
        "both models, computed delta-AUC on each, and reported the 2.5 and "
        "97.5 percentile of the differences plus the empirical two-sided p "
        "value. The enhanced model is above the screening model on AUC with "
        "95% confidence on this holdout.")


def build_utility(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Clinical utility and safety")
    add_title(slide, "Actionability matters more than a prettier metric.", size=30)

    # Left figure: DCA.
    add_image_centered(slide, FIG / "dca_combined.png",
                       PAD, Inches(2.20), Inches(6.40), Inches(4.30))
    add_caption(slide, PAD, Inches(6.50), Inches(6.40),
                "Decision-curve analysis: both selected thresholds sit inside useful "
                "net-benefit ranges.")

    # Right column: 2 cards + subgroup figure.
    add_card_content(slide, Inches(7.10), Inches(2.20),
                     Inches(5.80), Inches(1.20),
                     "Decision curve analysis",
                     "Both selected thresholds sit inside useful net-benefit ranges "
                     "against treat-all and treat-none baselines.",
                     accent=TEAL, accent_side="left", body_size=12, title_size=14.5)
    add_card_content(slide, Inches(7.10), Inches(3.55),
                     Inches(5.80), Inches(1.20),
                     "Subgroup honesty",
                     "Screening has a 25-point BMI recall gap, lowest in normal-BMI "
                     "patients. The app warns not to rule out lean PCOS solely from low "
                     "probability.",
                     accent=CORAL, accent_side="left", body_size=12, title_size=14.5)
    add_image_centered(slide, FIG / "subgroup_recall_by_bmi.png",
                       Inches(7.10), Inches(4.85), Inches(5.80), Inches(1.95))

    add_footer(slide, n, total, "7:05 – 7:45")
    set_notes(slide,
        "Decision curve analysis asks whether a model helps at clinically "
        "plausible thresholds — ours do. We also found an honest fairness "
        "risk: recall is lower in the normal-BMI group for screening. Rather "
        "than hide that, we added a warning to the prototype.")


def build_demo5(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Demo 5  ·  notebook 12  ·  appendix", color=CORAL)
    add_title(slide,
              "Split-conformal gives every patient a distribution-free coverage guarantee.",
              size=22, y=Inches(0.78), h=Inches(0.95))

    code_lines = [
        line([("c", "# 1. Split training into proper-train + calibration slice.")]),
        line([("fg", "X_pt, X_cal, y_pt, y_cal = "),
              ("f", "train_test_split"), ("fg", "(")]),
        line([("fg", "    X_train, y_train, test_size="), ("n", "0.30"), ("fg", ",")]),
        line([("fg", "    stratify=y_train, random_state="), ("n", "42"), ("fg", ",")]),
        line([("fg", ")")]),
        "",
        line([("c", "# 2. Refit a clone on PROPER-TRAIN only.")]),
        line([("fg", "model = "), ("f", "clone"),
              ("fg", "(enhanced_artifact["), ("s", '"model"'), ("fg", "])")]),
        line([("fg", "model."), ("f", "fit"), ("fg", "(X_pt, y_pt)")]),
        "",
        line([("c", "# 3. Nonconformity scores: s = 1 − p_model(y_true | x).")]),
        line([("fg", "p_cal  = model."), ("f", "predict_proba"), ("fg", "(X_cal)")]),
        line([("fg", "scores = "), ("n", "1"), ("fg", " - p_cal[np."),
              ("f", "arange"), ("fg", "("), ("f", "len"),
              ("fg", "(y_cal)), y_cal.values]")]),
        "",
        line([("c", "# 4. q-hat at ceil((n_cal+1)(1-alpha)) / n_cal.")]),
        line([("fg", "alpha = "), ("n", "0.10")]),
        line([("fg", "n_cal = "), ("f", "len"), ("fg", "(scores)")]),
        line([("fg", "q_hat = np."), ("f", "quantile"), ("fg", "(scores,")]),
        line([("fg", "                   np."), ("f", "ceil"),
              ("fg", "((n_cal + "), ("n", "1"), ("fg", ") * ("),
              ("n", "1"), ("fg", " - alpha)) / n_cal)")]),
        "",
        line([("c", "# 5. Prediction set: { y : 1 − p(y|x) <= q_hat }.")]),
        line([("fg", "p_te = model."), ("f", "predict_proba"), ("fg", "(X_test)")]),
        line([("fg", "sets = ("), ("n", "1"), ("fg", " - p_te) <= q_hat")]),
    ]
    code_card(slide, PAD, Inches(1.85), Inches(6.30), Inches(4.55),
              "notebook 12  —  split-conformal recipe", code_lines)

    out_lines = [
        line([("fg", "{")]),
        line([("fg", "  "), ("mu", '"target_coverage"'),
              ("fg", ":     0.900,")]),
        line([("fg", "  "), ("mu", '"q_hat"'), ("fg", ":               "),
              ("ok", "0.486"), ("fg", ",")]),
        line([("fg", "  "), ("mu", '"empirical_coverage"'),
              ("fg", ":  "), ("ok", "0.912"), ("fg", ",")]),
        line([("fg", "  "), ("mu", '"mean_set_size"'),
              ("fg", ":       0.993,")]),
        line([("fg", "  "), ("mu", '"singleton_fraction"'),
              ("fg", ":  "), ("ok", "0.993"), ("fg", ",")]),
        line([("fg", "  "), ("mu", '"doubleton_fraction"'),
              ("fg", ":  0.000,")]),
        line([("fg", "  "), ("mu", '"empty_fraction"'),
              ("fg", ":      0.007,")]),
        line([("fg", "  "), ("mu", '"worked_examples"'), ("fg", ": [")]),
        line([("fg", "    {"), ("mu", '"true"'), ("fg", ": 1, "),
              ("mu", '"p_pcos"'), ("fg", ": 0.903,")]),
        line([("fg", "     "), ("mu", '"set"'), ("fg", ": ["),
              ("ok", "1"), ("fg", "],   "), ("mu", '"covers"'),
              ("fg", ": "), ("ok", "true"), ("fg", "},")]),
        line([("fg", "    {"), ("mu", '"true"'), ("fg", ": 0, "),
              ("mu", '"p_pcos"'), ("fg", ": 0.271,")]),
        line([("fg", "     "), ("mu", '"set"'), ("fg", ": ["),
              ("ok", "0"), ("fg", "],   "), ("mu", '"covers"'),
              ("fg", ": "), ("ok", "true"), ("fg", "}")]),
        line([("fg", "  ]")]),
        line([("fg", "}")]),
    ]
    output_card(slide, Inches(7.10), Inches(1.85), Inches(5.75), Inches(4.55),
                "Output  —  outputs/metrics/conformal_coverage.json",
                out_lines)

    takeaway(slide, PAD, Inches(6.55), CONTENT_W, Inches(0.70),
        "Target coverage 0.90, empirical 0.912 on the held-out set. The 0.7% empty sets are "
        "honest abstentions — patients the model refuses to call confidently.")

    add_footer(slide, n, total, "appendix / Q&A")
    set_notes(slide,
        "Split-conformal gives each patient a prediction set with a "
        "distribution-free coverage guarantee. We hold out 30% of the "
        "training set as a calibration slice, compute one minus the model "
        "probability of the true label on it, and pick q-hat at the "
        "appropriate quantile. Targeted 90% coverage; observed 91.2% on the "
        "holdout. Empty sets count as abstentions.")


def build_explainability(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Explainability at the patient level")
    add_title(slide, "Every prediction needs a reason clinicians can inspect.",
              size=30)

    cards = [
        ("Patient-level SHAP",
         "The demo surfaces top log-odds contributors for the individual "
         "patient, not just global feature rankings.", TEAL),
        ("Checklist safety net",
         "The app shows which clinical evidence is present, which tests are "
         "missing, and when metabolic or lean-PCOS follow-up is needed.",
         CORAL),
    ]
    cx = PAD
    cy = Inches(2.30)
    cw = Inches(5.60)
    ch = Inches(1.85)
    for t, body, color in cards:
        add_card_content(slide, cx, cy, cw, ch, t, body, accent=color,
                         accent_side="left", body_size=13, title_size=15.5)
        cy += ch + Inches(0.20)

    add_image_centered(slide, FIG / "enhanced_shap_positive_case.png",
                       Inches(7.05), Inches(2.30), Inches(5.80), Inches(4.30))
    add_caption(slide, Inches(7.05), Inches(6.65), Inches(5.80),
                "Held-out positive case: follicle counts dominate the enhanced-model "
                "explanation.")

    add_footer(slide, n, total, "7:45 – 8:25")
    set_notes(slide,
        "Explainability is not an afterthought. The clinician sees the "
        "probability, the action threshold, the risk tier, and the top SHAP "
        "contributors. In the enhanced model, follicle counts rise to the top, "
        "which is clinically sensible. The checklist translates that into "
        "next-step reasoning.")


def build_demo4(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Demo 4  ·  notebook 05  ·  appendix", color=CORAL)
    add_title(slide,
              "SHAP attributes every prediction to the patient's actual feature values.",
              size=22, y=Inches(0.78), h=Inches(0.95))

    code_lines = [
        line([("k", "def "), ("f", "transform_through_preprocessing"),
              ("fg", "(model, X_row):")]),
        line([("fg", "    X = X_row.copy()")]),
        line([("k", "    for"), ("fg", " _, transformer "),
              ("k", "in"), ("fg", " model.steps[:-"),
              ("n", "1"), ("fg", "]:")]),
        line([("fg", "        X = transformer."), ("f", "transform"), ("fg", "(X)")]),
        line([("k", "    return"), ("fg", " X")]),
        "",
        line([("c", "# Reproduce the train/test split, pull positive holdout case.")]),
        line([("fg", "_, X_te, _, y_te = "),
              ("f", "train_test_split"), ("fg", "(")]),
        line([("fg", "    pcos[feats], pcos["), ("s", '"pcos_y_n"'), ("fg", "],")]),
        line([("fg", "    test_size="), ("n", "0.25"),
              ("fg", ", stratify=pcos["), ("s", '"pcos_y_n"'), ("fg", "],")]),
        line([("fg", "    random_state="), ("n", "42"), ("fg", ",")]),
        line([("fg", ")")]),
        line([("fg", "positive = pcos.loc[X_te[y_te == "), ("n", "1"),
              ("fg", "].index["), ("n", "0"), ("fg", "]]."),
              ("f", "to_dict"), ("fg", "()")]),
        "",
        line([("c", "# TreeExplainer is fast and exact for the RF.")]),
        line([("fg", "estimator   = artifact["), ("s", '"model"'),
              ("fg", "].named_steps["), ("s", '"model"'), ("fg", "]")]),
        line([("fg", "X_t         = "),
              ("f", "transform_through_preprocessing"), ("fg", "(artifact["),
              ("s", '"model"'), ("fg", "],")]),
        line([("fg", "                                            pd."),
              ("f", "DataFrame"), ("fg", "([positive])[feats])")]),
        line([("fg", "explainer   = shap."), ("f", "TreeExplainer"),
              ("fg", "(estimator)")]),
        line([("fg", "shap_values = explainer."), ("f", "shap_values"),
              ("fg", "(X_t)["), ("n", "0"), ("fg", ", :, "),
              ("n", "1"), ("fg", "]")]),
    ]
    code_card(slide, PAD, Inches(1.85), Inches(6.30), Inches(4.55),
              "notebook 05  —  shap_for_artifact (enhanced model, holdout positive)",
              code_lines)

    out_lines = [
        line([("mu", "Top SHAP drivers, known PCOS holdout case")]),
        line([("mu", "feature                value   log-odds Δ")]),
        line([("fg", "follicle_no_r          15.0     "),
              ("hi", "+0.213"), ("fg", "  ↑ PCOS")]),
        line([("fg", "follicle_no_l          13.0     "),
              ("hi", "+0.166"), ("fg", "  ↑ PCOS")]),
        line([("fg", "skin_darkening_y_n      0.0     "),
              ("ok", "−0.040"), ("fg", "  ↓ PCOS")]),
        line([("fg", "weight_gain_y_n         0.0     "),
              ("ok", "−0.036"), ("fg", "  ↓ PCOS")]),
        line([("fg", "amh_ng_ml              16.4     "),
              ("hi", "+0.034"), ("fg", "  ↑ PCOS")]),
        line([("fg", "cycle_irregular_flag    1.0     "),
              ("hi", "+0.028"), ("fg", "  ↑ PCOS")]),
        line([("fg", "bmi                    27.3     "),
              ("hi", "+0.022"), ("fg", "  ↑ PCOS")]),
        line([("fg", "fsh_lh                  0.39    "),
              ("hi", "+0.018"), ("fg", "  ↑ PCOS")]),
    ]
    output_card(slide, Inches(7.10), Inches(1.85), Inches(5.75), Inches(4.55),
                "Output  —  outputs/metrics/demo_positive_enhanced_shap.csv",
                out_lines)

    takeaway(slide, PAD, Inches(6.55), CONTENT_W, Inches(0.70),
        "A clinician opens the patient record and sees exactly which findings pushed the "
        "prediction up or down. Demo cases come from held-out rows the model never saw in training.")

    add_footer(slide, n, total, "appendix / Q&A")
    set_notes(slide,
        "For one held-out positive case, SHAP tells us why the model said "
        "PCOS: high follicle counts on both ovaries do most of the work, with "
        "AMH and cycle irregularity contributing, and the absence of skin "
        "darkening and weight gain pulling slightly the other way.")


def build_benchmarks(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "External benchmarks")
    add_title(slide, "We benchmarked against both modern ML and clinical rules.",
              size=30)

    # Left: TabPFN figure.
    add_image_centered(slide, FIG / "tabpfn_vs_rf_auc.png",
                       PAD, Inches(2.20), Inches(6.40), Inches(4.30))
    add_caption(slide, PAD, Inches(6.50), Inches(6.40),
                "TabPFN-v2 (foundation model) vs the random-forest pipelines.")

    # Right: 2 cards + Rotterdam figure.
    add_card_content(slide, Inches(7.10), Inches(2.20),
                     Inches(5.80), Inches(1.30),
                     "TabPFN-v2 benchmark",
                     "AUC 0.905 screening / 0.962 enhanced; trades away recall at our "
                     "clinical threshold and is slower on CPU.",
                     accent=TEAL, accent_side="left", body_size=12, title_size=14.5)
    add_card_content(slide, Inches(7.10), Inches(3.65),
                     Inches(5.80), Inches(1.20),
                     "Rotterdam comparison",
                     "Enhanced ML adds 9.1 sensitivity points with a 1.1 specificity-point "
                     "cost versus a hand-coded 2-of-3 rule.",
                     accent=CORAL, accent_side="left", body_size=12, title_size=14.5)
    add_image_centered(slide, FIG / "rotterdam_vs_ml.png",
                       Inches(7.10), Inches(4.95), Inches(5.80), Inches(1.85))

    add_footer(slide, n, total, "8:25 – 9:00")
    set_notes(slide,
        "We did not stop at comparing random forest to logistic regression. "
        "TabPFN slightly improves AUC, but the random forest keeps better "
        "recall and F2 at the selected action threshold and is easier to "
        "deploy. Against Rotterdam, enhanced ML catches more positives with "
        "only a small specificity cost.")


def build_differential(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Differential diagnosis prompt")
    add_title(slide, "PCOS Pathfinder avoids tunnel vision.", size=30)

    tb, tf = add_text_box(slide, PAD, Inches(2.10), Inches(7.40), Inches(1.0),
                          margin=Inches(0))
    add_paragraph(tf, "The endometriosis module is intentionally modest: a "
                  "synthetic-data prompt that says \"consider broader workup,\" "
                  "not \"diagnose endometriosis.\"",
                  size=14, color=MUTED, first=True)

    # 3 metric cards.
    metrics = [
        ("0.660", "Overlap AUC",       TEAL),
        ("6",     "Prompt features",   CORAL),
        ("0",     "Diagnostic overclaim", GOLD),
    ]
    card_w = Inches(2.30)
    card_h = Inches(1.60)
    x = PAD
    y = Inches(3.50)
    for value, label, color in metrics:
        add_metric_card(slide, x, y, card_w, card_h, value, label, accent=color)
        x += card_w + Inches(0.20)

    # Right: confusion matrix.
    add_image_centered(slide, FIG / "endometriosis_confusion_matrix.png",
                       Inches(7.95), Inches(2.20), Inches(4.90), Inches(4.40))

    add_footer(slide, n, total, "9:00 – 9:30")
    set_notes(slide,
        "The differential module is framed as a prompt. Because the "
        "endometriosis dataset is synthetic, we make no clinical diagnostic "
        "claim. But when chronic pain, infertility, hormonal abnormality and "
        "menstrual irregularity cluster together, the system reminds the "
        "clinician not to stop thinking after the PCOS risk score.")


def build_implementation(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Implementation path")
    add_title(slide, "Submission-ready implementation, not a concept sketch.",
              size=30)

    cards = [
        ("Prototype",
         "Streamlit app loads the same joblib artifacts generated by the notebooks.",
         TEAL),
        ("Reproducibility",
         "13 executed notebooks + 15-page technical report PDF · deterministic "
         "notebook builders · pinned dependencies.",
         CORAL),
        ("Artifacts",
         "Metrics JSON/CSV, figures, model card, cleaned data, saved models.",
         GOLD),
        ("Packaging",
         "README, submission checklist, ignored large source archives, documented limitations.",
         LAV),
    ]
    card_w = (CONTENT_W - Inches(0.30) * 3) / 4
    card_h = Inches(1.95)
    y = Inches(2.20)
    x = PAD
    for t, body, color in cards:
        add_card_content(slide, x, y, card_w, card_h, t, body, accent=color,
                         body_size=11.5, title_size=14)
        x += card_w + Inches(0.30)

    headers = ["Deployment step", "What changes before real use", "Acceptance gate"]
    rows = [
        ["Internal validation", "Repeat notebook pipeline on partner data",
         "Recall and calibration remain stable"],
        ["Clinical pilot",     "Run as silent decision support beside usual care",
         "Measures earlier referral without excess harm"],
        ["Prospective study",  "Collect outcomes and subgroup performance",
         "External validation before deployment"],
    ]
    _add_table(slide, PAD, Inches(4.55), CONTENT_W, headers, rows,
               col_widths=[Inches(3.30), Inches(4.80), Inches(4.10)],
               row_h=Inches(0.50))

    add_footer(slide, n, total, "10:40 – 11:20")
    set_notes(slide,
        "This is feasible because it is already implemented. The Streamlit "
        "app loads saved artifacts, the notebooks reproduce the metrics, and "
        "the outputs are packaged for review. But we do not pretend this is "
        "ready for clinical deployment. The next step is external validation, "
        "then a silent pilot, then prospective evaluation.")


def build_impact(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Patient and system impact")
    add_title(slide, "Impact comes from earlier, clearer next steps.", size=30)

    bullets = [
        ("Primary care",
         "Identify patients who need PCOS assessment before years of delay."),
        ("Low-resource settings",
         "The screening model uses symptoms and basic vitals only."),
        ("Patient safety",
         "Risk tiers, SHAP, calibration, conformal abstention, and explicit caveats."),
        ("Equity honesty",
         "Subgroup gaps are reported, not hidden."),
    ]
    tb, tf = add_text_box(slide, PAD, Inches(2.20), Inches(7.50), Inches(4.80),
                          margin=Inches(0))
    for i, (h, body) in enumerate(bullets):
        add_runs_paragraph(tf, [
            {"text": "•   ", "size": 17, "color": TEAL, "bold": True},
            {"text": f"{h}: ", "size": 16, "color": NAVY, "bold": True},
            {"text": body, "size": 15.5, "color": MUTED},
        ], first=(i == 0), space_after=14)

    # Right metrics.
    metrics = [
        ("13",    "Features for frontline screening",            TEAL),
        ("0.912", "Conformal coverage on holdout",               CORAL),
        ("20",    "Curated PCOS genes in sampled feature list",  GOLD),
    ]
    cy = Inches(2.20)
    cx = Inches(8.30)
    cw = Inches(4.50)
    ch = Inches(1.45)
    for value, label, color in metrics:
        card = add_card(slide, cx, cy, cw, ch, accent=color, accent_side="left")
        tb, tf = add_text_box(slide, cx, cy + Inches(0.15),
                              cw, ch - Inches(0.30), margin=Inches(0.10))
        add_paragraph(tf, value, size=36, color=color, bold=True,
                      align="center", space_after=4, first=True)
        add_paragraph(tf, label, size=11.5, color=MUTED, align="center")
        cy += ch + Inches(0.18)

    add_footer(slide, n, total, "11:20 – 12:00")
    set_notes(slide,
        "The public-health value is earlier triage. The low-resource value is "
        "that the first model does not require labs or ultrasound. The safety "
        "value is that the system explains itself and states what it does not "
        "know. And the equity value is honesty: we found a subgroup risk and "
        "built a warning around it, rather than presenting a polished but "
        "unsafe story.")


# ---------------------------------------------------------------------------
# App-backup slides (static replicas of the Streamlit prototype)
# ---------------------------------------------------------------------------
def _app_intake_block(slide, x, y, w, h, header, intake_rows, checklist_rows):
    card = add_card(slide, x, y, w, h)
    tb, tf = add_text_box(slide, x + Inches(0.20), y + Inches(0.18),
                          w - Inches(0.40), h - Inches(0.30))
    add_paragraph(tf, header, size=12, color=TEAL, bold=True,
                  space_after=6, first=True)
    for label, value, tone in intake_rows:
        color = INK if tone is None else (CORAL if tone == "pos" else GREEN)
        add_runs_paragraph(tf, [
            {"text": label, "size": 11, "color": MUTED},
            {"text": value, "size": 11, "color": color, "bold": True},
        ], space_after=2)
        # Right-align numeric — done by tab-style spacer in label.
    add_paragraph(tf, "Diagnostic completeness", size=12, color=TEAL,
                  bold=True, space_before=8, space_after=4)
    for tick, body in checklist_rows:
        color = GREEN if tick else MUTED
        glyph = "✓ " if tick else "—  "
        add_runs_paragraph(tf, [
            {"text": glyph, "size": 12, "color": color, "bold": True},
            {"text": body, "size": 11, "color": INK},
        ], space_after=3)


def _app_assessment_block(slide, x, y, w, h, header, tiers, shap_rows, note):
    card = add_card(slide, x, y, w, h)
    tb, tf = add_text_box(slide, x + Inches(0.20), y + Inches(0.18),
                          w - Inches(0.40), Inches(0.35))
    add_paragraph(tf, header, size=12, color=TEAL, bold=True, first=True)

    # Tier result cards.
    ty = y + Inches(0.55)
    th = Inches(0.85)
    for i, (name, tag, prob, thr, tier_color, tier_label) in enumerate(tiers):
        row = add_card(slide, x + Inches(0.20), ty + i * (th + Inches(0.10)),
                       w - Inches(0.40), th, fill=RGBColor(0xf8, 0xfa, 0xfc),
                       accent=tier_color, accent_side="left")
        tb, tf = add_text_box(slide, x + Inches(0.35),
                              ty + i * (th + Inches(0.10)),
                              w - Inches(0.60), th, anchor="middle",
                              margin=Inches(0.0))
        add_runs_paragraph(tf, [
            {"text": name, "size": 12, "color": NAVY, "bold": True},
            {"text": f"   {tag}", "size": 10, "color": MUTED},
        ], first=True, space_after=2)
        add_runs_paragraph(tf, [
            {"text": "Probability ", "size": 10, "color": MUTED},
            {"text": prob, "size": 13, "color": tier_color, "bold": True},
            {"text": f"     Threshold {thr}     ", "size": 10, "color": MUTED},
            {"text": tier_label, "size": 11, "color": tier_color, "bold": True},
        ])

    # SHAP table.
    sy = ty + 2 * (th + Inches(0.10)) + Inches(0.15)
    tb, tf = add_text_box(slide, x + Inches(0.20), sy,
                          w - Inches(0.40), Inches(0.30))
    add_paragraph(tf, "Top SHAP drivers (enhanced model)", size=12,
                  color=TEAL, bold=True, first=True)

    table_y = sy + Inches(0.32)
    headers = ["Feature", "Value", "Log-odds", "Direction"]
    _add_table(slide, x + Inches(0.20), table_y, w - Inches(0.40),
               headers, shap_rows,
               col_widths=[Inches(1.80), Inches(0.80), Inches(1.00), Inches(1.20)],
               row_h=Inches(0.30), header_h=Inches(0.32))

    # Note at the bottom.
    note_y = y + h - Inches(0.42)
    tb, tf = add_text_box(slide, x + Inches(0.20), note_y,
                          w - Inches(0.40), Inches(0.32))
    add_paragraph(tf, note, size=10, color=MUTED, italic=True, first=True)


def build_app_high(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "App walkthrough  ·  high-risk case", color=CORAL)
    add_title(slide, "What the clinician sees — held-out PCOS-positive patient.",
              size=24, y=Inches(0.78), h=Inches(0.95))

    intake = [
        ("Age              ",  "29",            None),
        ("BMI              ",  "29.3",          "pos"),
        ("Cycle            ",  "Irregular",     "pos"),
        ("Bleeding days    ",  "3",             None),
        ("Weight gain      ",  "Yes",           "pos"),
        ("Hair growth      ",  "No",            "neg"),
        ("Skin darkening   ",  "Yes",           "pos"),
        ("Hair loss        ",  "Yes",           "pos"),
        ("Pimples          ",  "Yes",           "pos"),
        ("RBS              ",  "92 mg/dl",      None),
        ("AMH              ",  "6.0 ng/mL",     "pos"),
        ("FSH/LH           ",  "0.90",          None),
        ("Follicles L/R    ",  "10 / 9",        "pos"),
        ("Endometrium      ",  "8 mm",          None),
    ]
    checklist = [
        (True, "Ovulatory dysfunction: irregular cycle pattern present"),
        (True, "Clinical hyperandrogenism: skin darkening + hair loss + acne"),
        (True, "Polycystic morphology: follicle counts 10 / 9"),
        (True, "Metabolic follow-up: consider (BMI 29.3, overweight)"),
    ]
    _app_intake_block(slide, PAD, Inches(1.85), Inches(5.80), Inches(5.40),
                      "Patient intake (test-set row 395)", intake, checklist)

    tiers = [
        ("Frontline screening", "primary care / telehealth, no labs",
         "88.6%", "0.29", CORAL, "HIGH RISK"),
        ("Enhanced (labs + ultrasound)", "diagnostic support",
         "90.7%", "0.38", CORAL, "HIGH RISK"),
    ]
    shap_rows = [
        ["follicle_no_r",  "9",  "+0.092", "↑ PCOS"],
        ["follicle_no_l",  "10", "+0.090", "↑ PCOS"],
        ["skin_darkening", "1",  "+0.069", "↑ PCOS"],
        ["weight_gain",    "1",  "+0.066", "↑ PCOS"],
        ["hair_growth",    "0",  "−0.031", "↓ PCOS"],
    ]
    _app_assessment_block(slide, Inches(6.85), Inches(1.85),
                          Inches(6.00), Inches(5.40),
                          "Tiered assessment", tiers, shap_rows,
                          "Pre-rendered from outputs/metrics/demo_*.csv. "
                          "Equivalent to `streamlit run src/app.py`.")

    add_footer(slide, n, total, "9:30 – 10:05")
    set_notes(slide,
        "Static app backup. The clinician fills the intake; both tiered "
        "scores land in High-risk; the top SHAP drivers are led by ovarian "
        "morphology; the Rotterdam-style checklist has all four boxes ticked. "
        "The whole flow is below the fold of a single laptop screen.")


def build_app_low(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "App walkthrough  ·  low-risk case", color=GREEN)
    add_title(slide, "What the clinician sees — held-out non-PCOS patient.",
              size=24, y=Inches(0.78), h=Inches(0.95))

    intake = [
        ("Age              ", "32",            None),
        ("BMI              ", "24.2",          None),
        ("Cycle            ", "Regular",       None),
        ("Bleeding days    ", "6",             None),
        ("Weight gain      ", "Yes",           "pos"),
        ("Hair growth      ", "No",            "neg"),
        ("Skin darkening   ", "No",            "neg"),
        ("Hair loss        ", "No",            "neg"),
        ("Pimples          ", "Yes",           "pos"),
        ("RBS              ", "100 mg/dl",     None),
        ("AMH              ", "11.4 ng/mL",    "pos"),
        ("FSH/LH           ", "3.03",          None),
        ("Follicles L/R    ", "2 / 3",         None),
        ("Endometrium      ", "8 mm",          None),
    ]
    checklist = [
        (False, "Ovulatory dysfunction: regular cycle pattern"),
        (False, "Clinical hyperandrogenism: only acne reported"),
        (False, "Polycystic morphology: follicle counts 2 / 3 (low)"),
        (False, "Metabolic follow-up: no immediate prompt (BMI 24.2)"),
    ]
    _app_intake_block(slide, PAD, Inches(1.85), Inches(5.80), Inches(5.40),
                      "Patient intake (test-set row 383)", intake, checklist)

    tiers = [
        ("Frontline screening", "primary care / telehealth, no labs",
         "20.8%", "0.29", GREEN, "LOW RISK"),
        ("Enhanced (labs + ultrasound)", "diagnostic support",
         "21.4%", "0.38", GREEN, "LOW RISK"),
    ]
    shap_rows = [
        ["follicle_no_r",  "3",   "−0.157", "↓ PCOS"],
        ["follicle_no_l",  "2",   "−0.117", "↓ PCOS"],
        ["amh_ng_ml",      "11.4", "+0.062", "↑ PCOS"],
        ["weight_gain",    "1",   "+0.061", "↑ PCOS"],
        ["skin_darkening", "0",   "−0.036", "↓ PCOS"],
    ]
    _app_assessment_block(slide, Inches(6.85), Inches(1.85),
                          Inches(6.00), Inches(5.40),
                          "Tiered assessment", tiers, shap_rows,
                          "Note the elevated AMH (11.4) pushes up, but low follicle counts "
                          "(2 / 3) dominate — model correctly says low risk.")

    add_footer(slide, n, total, "10:05 – 10:40")
    set_notes(slide,
        "The second backup case: AMH is elevated at 11.4, which a clinician "
        "might initially worry about. But ovarian morphology shows low "
        "follicle counts of 2 and 3, and the cycle is regular. The model "
        "integrates these signals correctly and returns low risk on both "
        "tiers.")


# ---------------------------------------------------------------------------
# Combined app-walkthrough slide (high-risk + low-risk side by side)
# ---------------------------------------------------------------------------
def _app_compact_panel(slide, x, y, w, h, *, case_label, case_color,
                       intake_rows, tiers, shap_rows, note):
    """One half of the combined app slide — compact patient card."""
    add_card(slide, x, y, w, h, accent=case_color, accent_side="left")

    inner_x = x + Inches(0.22)
    inner_w = w - Inches(0.44)

    # Header.
    tb, tf = add_text_box(slide, inner_x, y + Inches(0.16),
                          inner_w, Inches(0.36))
    add_paragraph(tf, case_label, size=13, color=case_color, bold=True,
                  first=True)

    # Intake summary — 5 most relevant rows.
    intake_y = y + Inches(0.58)
    tb, tf = add_text_box(slide, inner_x, intake_y, inner_w, Inches(1.50))
    add_paragraph(tf, "Patient intake (key features)", size=11, color=TEAL,
                  bold=True, space_after=4, first=True)
    for label, value, tone in intake_rows:
        color = INK if tone is None else (CORAL if tone == "pos" else GREEN)
        add_runs_paragraph(tf, [
            {"text": label, "size": 10.5, "color": MUTED},
            {"text": value, "size": 10.5, "color": color, "bold": True},
        ], space_after=2)

    # Tiered scores — 2 compact rows.
    tiered_y = intake_y + Inches(1.85)
    tb, tf = add_text_box(slide, inner_x, tiered_y, inner_w, Inches(0.30))
    add_paragraph(tf, "Tiered assessment", size=11, color=TEAL, bold=True,
                  first=True)
    ty = tiered_y + Inches(0.32)
    th = Inches(0.42)
    for i, (name, prob, thr, tier_color, tier_label) in enumerate(tiers):
        ry = ty + i * (th + Inches(0.06))
        add_card(slide, inner_x, ry, inner_w, th,
                 fill=RGBColor(0xf8, 0xfa, 0xfc),
                 accent=tier_color, accent_side="left")
        tb, tf = add_text_box(slide, inner_x + Inches(0.18), ry,
                              inner_w - Inches(0.25), th, anchor="middle",
                              margin=Inches(0.0))
        add_runs_paragraph(tf, [
            {"text": name, "size": 10.5, "color": NAVY, "bold": True},
            {"text": "   p=", "size": 10, "color": MUTED},
            {"text": prob, "size": 11.5, "color": tier_color, "bold": True},
            {"text": f"  thr {thr}   ", "size": 10, "color": MUTED},
            {"text": tier_label, "size": 10.5, "color": tier_color, "bold": True},
        ], first=True)

    # Top SHAP drivers — render as styled text rows (not a table, since
    # PowerPoint can choke on multiple tables-inside-cards on one slide).
    shap_y = ty + 2 * (th + Inches(0.06)) + Inches(0.15)
    tb, tf = add_text_box(slide, inner_x, shap_y, inner_w, Inches(1.40))
    add_paragraph(tf, "Top SHAP drivers (enhanced)", size=11, color=TEAL,
                  bold=True, space_after=4, first=True)
    for feature, value, logodds, direction in shap_rows:
        dir_color = CORAL if direction == "↑" else GREEN
        add_runs_paragraph(tf, [
            {"text": f"{feature:<18}", "font": "Consolas", "size": 10.5,
             "color": INK},
            {"text": f"{value:>6}   ", "font": "Consolas", "size": 10.5,
             "color": MUTED},
            {"text": f"{logodds:>7}  ", "font": "Consolas", "size": 10.5,
             "color": dir_color, "bold": True},
            {"text": direction, "font": "Consolas", "size": 11,
             "color": dir_color, "bold": True},
        ], space_after=2)

    # Footer note.
    note_y = y + h - Inches(0.42)
    tb, tf = add_text_box(slide, inner_x, note_y, inner_w, Inches(0.35))
    add_paragraph(tf, note, size=9.5, color=MUTED, italic=True, first=True)


def build_app_combined(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "App walkthrough  ·  two held-out cases")
    add_title(slide,
              "What the clinician sees — high-risk vs. low-risk patient.",
              size=24, y=Inches(0.78), h=Inches(0.95))

    col_w = (CONTENT_W - Inches(0.30)) / 2
    col_h = Inches(5.40)
    col_y = Inches(1.85)

    # High-risk case (left).
    _app_compact_panel(
        slide, PAD, col_y, col_w, col_h,
        case_label="High-risk  ·  test-set row 395  ·  PCOS-positive (true)",
        case_color=CORAL,
        intake_rows=[
            ("BMI              ", "29.3",       "pos"),
            ("Cycle            ", "Irregular",  "pos"),
            ("Skin darkening   ", "Yes",        "pos"),
            ("AMH              ", "6.0 ng/mL",  "pos"),
            ("Follicles L/R    ", "10 / 9",     "pos"),
        ],
        tiers=[
            ("Frontline screening",            "88.6%", "0.29", CORAL, "HIGH"),
            ("Enhanced (labs + ultrasound)",  "90.7%", "0.38", CORAL, "HIGH"),
        ],
        shap_rows=[
            ["follicle_no_r",  "9",  "+0.092", "↑"],
            ["follicle_no_l",  "10", "+0.090", "↑"],
            ["skin_darkening", "1",  "+0.069", "↑"],
        ],
        note="All four Rotterdam-style checklist boxes tick. SHAP led by ovarian morphology.",
    )

    # Low-risk case (right).
    _app_compact_panel(
        slide, PAD + col_w + Inches(0.30), col_y, col_w, col_h,
        case_label="Low-risk  ·  test-set row 383  ·  non-PCOS (true)",
        case_color=GREEN,
        intake_rows=[
            ("BMI              ", "24.2",       None),
            ("Cycle            ", "Regular",    None),
            ("Skin darkening   ", "No",         "neg"),
            ("AMH              ", "11.4 ng/mL", "pos"),
            ("Follicles L/R    ", "2 / 3",      None),
        ],
        tiers=[
            ("Frontline screening",           "20.8%", "0.29", GREEN, "LOW"),
            ("Enhanced (labs + ultrasound)", "21.4%", "0.38", GREEN, "LOW"),
        ],
        shap_rows=[
            ["follicle_no_r",  "3",    "−0.157", "↓"],
            ["follicle_no_l",  "2",    "−0.117", "↓"],
            ["amh_ng_ml",      "11.4", "+0.062", "↑"],
        ],
        note="Elevated AMH pushes up, but low follicle counts dominate — correctly Low Risk.",
    )

    add_footer(slide, n, total)
    set_notes(slide,
        "Two held-out app walkthroughs side by side. The high-risk patient: "
        "follicle counts ten and nine, irregular cycle, skin darkening — both "
        "tiers say High Risk, all four checklist boxes tick. The low-risk "
        "patient: elevated AMH at 11.4 might worry a clinician, but low "
        "follicle counts and regular cycle correctly drive a Low Risk score on "
        "both tiers. Pre-rendered from outputs/metrics/demo_*.csv.")


# ---------------------------------------------------------------------------
# Combined closing slide (implementation path + impact)
# ---------------------------------------------------------------------------
def build_close(prs, n, total):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, "Implementation path  ·  patient & system impact")
    add_title(slide,
              "From submission-ready to deployable — honestly.",
              size=28)

    # Top section: 4 impact themes as horizontal cards.
    themes = [
        ("Earlier triage",
         "Identify patients who need PCOS assessment before years of delay.",
         TEAL),
        ("Low-resource fit",
         "Frontline tier uses symptoms and basic vitals — no labs needed.",
         CORAL),
        ("Patient safety",
         "SHAP, calibration, conformal abstention, and explicit caveats.",
         GOLD),
        ("Equity honesty",
         "Subgroup gaps are measured, surfaced in the app, and flagged in the model card.",
         LAV),
    ]
    card_w = (CONTENT_W - Inches(0.30) * 3) / 4
    card_h = Inches(1.85)
    y = Inches(2.05)
    x = PAD
    for t, body, color in themes:
        add_card_content(slide, x, y, card_w, card_h, t, body, accent=color,
                         body_size=11.5, title_size=14)
        x += card_w + Inches(0.30)

    # Bottom section: deployment-path table.
    headers = ["Deployment step", "What changes before real use", "Acceptance gate"]
    rows = [
        ["Internal validation",
         "Repeat notebook pipeline on partner data",
         "Recall and calibration remain stable"],
        ["Clinical pilot",
         "Run as silent decision support beside usual care",
         "Measures earlier referral without excess harm"],
        ["Prospective study",
         "Collect outcomes and subgroup performance",
         "External validation before deployment"],
    ]
    tbl_y = y + card_h + Inches(0.45)
    _add_table(slide, PAD, tbl_y, CONTENT_W, headers, rows,
               col_widths=[Inches(3.30), Inches(4.80), Inches(4.10)],
               row_h=Inches(0.50))

    add_footer(slide, n, total)
    set_notes(slide,
        "The deployment path is honest: internal validation, silent pilot, "
        "prospective study — external validation before any clinical claim. "
        "The impact is earlier triage for delayed patients, a screening model "
        "that needs no labs for low-resource settings, transparent SHAP and "
        "conformal abstention for safety, and subgroup gaps that we measured "
        "and surfaced rather than hid. Thank you.")


# ---------------------------------------------------------------------------
# Biology appendix slides
# ---------------------------------------------------------------------------
def _build_bio_slide(prs, n, total, eyebrow, title, image_name, summary, notes,
                     time_label="1:00 – 1:35"):
    slide = add_blank_slide(prs)
    add_accent_strip(slide)
    add_eyebrow(slide, eyebrow, color=LAV)
    add_title(slide, title, size=26)

    # Left: image.
    add_image_centered(slide, BIO_FIG / image_name,
                       PAD, Inches(2.10), Inches(8.20), Inches(4.85))

    # Right: summary card.
    add_card(slide, Inches(9.00), Inches(2.10),
             Inches(3.85), Inches(4.85), accent=LAV, accent_side="left")
    tb, tf = add_text_box(slide, Inches(9.20), Inches(2.30),
                          Inches(3.55), Inches(4.55))
    add_paragraph(tf, "Why this matters for the model", size=14, color=LAV,
                  bold=True, space_after=8, first=True)
    for bullet in summary:
        add_runs_paragraph(tf, [
            {"text": "•  ", "size": 12, "color": LAV, "bold": True},
            {"text": bullet, "size": 12, "color": INK},
        ], space_after=8)

    add_footer(slide, n, total, time_label)
    set_notes(slide, notes)


def build_bio_hpo(prs, n, total):
    _build_bio_slide(prs, n, total,
        "Pathophysiology  ·  HPO axis",
        "The hyperandrogenism in PCOS is upstream of the ovary.",
        "hpo_axis.png",
        summary=[
            "PCOS disrupts GnRH pulse frequency, raising LH relative to FSH.",
            "Elevated LH drives theca-cell androgen production; relative FSH "
            "deficit reduces aromatization to estradiol.",
            "Clinically: hirsutism, acne, oligo-ovulation — features the screening "
            "tier uses as low-friction inputs.",
            "Why an enhanced tier matters: testosterone is the missing direct "
            "biomarker; AMH + LH:FSH context partially substitutes.",
        ],
        notes=(
            "PCOS hyperandrogenism is not a primary ovarian problem alone — the "
            "GnRH pulse generator runs fast, LH dominates relative to FSH, and "
            "theca cells overproduce androgens. The screening model picks up the "
            "downstream skin and cycle features without needing the lab; the "
            "enhanced tier adds AMH and follicle counts that reflect the "
            "follicular-arrest endpoint."),
        time_label="1:00 – 1:35")


def build_bio_insulin(prs, n, total):
    _build_bio_slide(prs, n, total,
        "Pathophysiology  ·  Metabolic axis",
        "Insulin resistance amplifies the androgen loop and drives long-term risk.",
        "insulin_androgen.png",
        summary=[
            "Peripheral insulin resistance raises insulin, which directly "
            "stimulates ovarian theca androgen output.",
            "Hyperinsulinemia also lowers hepatic SHBG, raising free testosterone.",
            "PCOS confers 3–7× relative risk of T2DM; CVD, GDM, NAFLD, and "
            "endometrial hyperplasia all elevated.",
            "Why our app warns on lean PCOS: metabolic risk is not purely BMI-driven; "
            "lean patients can still be insulin-resistant.",
        ],
        notes=(
            "The insulin–androgen loop is the metabolic engine of PCOS. The "
            "long-term risk multipliers — type-2 diabetes, cardiovascular events, "
            "NAFLD — justify treating PCOS as a metabolic-cardiovascular condition, "
            "not only a reproductive one. The app's lean-PCOS warning is grounded "
            "in this biology: low BMI does not rule out insulin resistance."),
        time_label="1:35 – 2:10")


def build_bio_ovary(prs, n, total):
    _build_bio_slide(prs, n, total,
        "Morphology  ·  Follicular arrest",
        "Polycystic morphology is follicular arrest, not multiple cysts.",
        "ovary_morphology.png",
        summary=[
            "Healthy ovary: 5–10 antral follicles, one dominant follicle matures "
            "and ovulates.",
            "PCOM: ≥ 12 follicles 2–9 mm arranged peripherally; dense stroma; "
            "no dominant follicle.",
            "Rotterdam 2003 ultrasound criterion: ≥12 follicles or volume ≥10 mL.",
            "Why follicle counts dominate enhanced-model SHAP: they are the "
            "most specific morphologic feature available in this dataset.",
        ],
        notes=(
            "Polycystic morphology is not multiple cysts — it is follicular "
            "arrest. Small antral follicles accumulate at the periphery because "
            "none receive sufficient FSH to become dominant. This is why "
            "follicle_no_l and follicle_no_r dominate enhanced-model SHAP for "
            "positive cases: they are the most specific morphologic signal in "
            "the dataset."),
        time_label="2:10 – 2:45")


def build_bio_phenotypes(prs, n, total):
    _build_bio_slide(prs, n, total,
        "Heterogeneity  ·  Rotterdam phenotypes",
        "PCOS is heterogeneous: 4 Rotterdam phenotypes with different presentations.",
        "phenotype_venn.png",
        summary=[
            "Rotterdam criteria: any 2 of hyperandrogenism, ovulatory dysfunction, "
            "polycystic morphology.",
            "Phenotype A (all three) ≈50%; B (H+O) ≈10%; C (H+P) ≈25%; "
            "D (O+P) ≈15% (Lizneva 2016).",
            "Phenotype B can lack ultrasound findings — screening tier catches "
            "it without follicle counts.",
            "Lean PCOS overlaps multiple phenotypes; the app's lean-PCOS warning "
            "addresses the screening recall gap we measured.",
        ],
        notes=(
            "PCOS is not one disease. Four phenotypes presenting through "
            "different combinations of hyperandrogenism, ovulatory dysfunction, "
            "and polycystic morphology — with different metabolic profiles. The "
            "screening tier exists in part to catch phenotypes B (H+O, no PCOM) "
            "without requiring ultrasound."),
        time_label="2:45 – 3:20")


def build_bio_singlecell(prs, n, total):
    _build_bio_slide(prs, n, total,
        "Molecular rationale  ·  Single-cell overlap",
        "Single-cell gene programs ground the diagnostic feature set in biology.",
        "single_cell_genes.png",
        summary=[
            "Curated PCOS gene programs span steroidogenesis, androgen signalling, "
            "insulin axis, follicular development, and GWAS loci.",
            "20 genes recovered from a sampled single-cell ovary feature list — "
            "biological rationale, not an expression analysis.",
            "Same axes (androgens, insulin, follicular development) that the "
            "tabular features capture downstream.",
            "Limitation: the diagnostic model uses tabular clinical features only — "
            "this slide grounds those features in the underlying biology.",
        ],
        notes=(
            "The single-cell layer is used for biological rationale only, not as "
            "expression evidence. The point is that the same biological axes that "
            "drive PCOS at the molecular level — steroidogenesis, androgen "
            "signalling, insulin, follicular development, GWAS loci — are the "
            "axes our clinical features capture downstream."),
        time_label="3:20 – 3:55")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Set default note master font.
    builders = [
        # Main flow — 18 slides, 9-minute pitch.
        # A — Problem & biology (5 slides directly serve C&S Validity, 30%).
        build_title,
        build_problem,
        build_bio_hpo,
        build_bio_insulin,
        build_bio_ovary,
        build_bio_phenotypes,
        build_bio_singlecell,
        # B — Solution design.
        build_workflow,
        build_data,
        build_method,
        # C — Results & rigor.
        build_results,
        build_rigor,
        build_utility,
        build_explainability,
        # D — Differentiation.
        build_benchmarks,
        build_differential,
        # E — Prototype & closing.
        build_app_combined,
        build_close,
        # Appendix — 5 live-code demos, available for Q&A.
        build_demo1,
        build_demo2,
        build_demo3,
        build_demo4,
        build_demo5,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(prs, i, total)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX}  ({total} slides)")


if __name__ == "__main__":
    main()
